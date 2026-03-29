"""
generate_report.py — APT 실거래 데이터 분석 PPTX 보고서 생성기
사용법: python scripts/generate_report.py --format pptx
"""
import argparse
import os
import platform
import glob
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── 경로 설정 ──
BASE_DIR = Path(__file__).parent.parent  # test/
DOCS_DIR = BASE_DIR / "docs"
FIGURES_DIR = DOCS_DIR / "figures"
OUTPUT_PPTX = DOCS_DIR / "report.pptx"

# ── 폰트 ──
FONT = 'AppleGothic' if platform.system() == 'Darwin' else 'Malgun Gothic'

# ── 색상 테마 (viz/STYLE_GUIDE.md 동기화) ──
CLR_TITLE = RGBColor(0x33, 0x33, 0x33)
CLR_BODY = RGBColor(0x55, 0x55, 0x55)
CLR_ACCENT = RGBColor(0x4A, 0x90, 0xD9)
CLR_ALERT = RGBColor(0xE7, 0x4C, 0x3C)
CLR_TBL_HEADER = RGBColor(0x2C, 0x3E, 0x50)
CLR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CLR_LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
CLR_COVER_BG = RGBColor(0x2C, 0x3E, 0x50)

# ── 프로젝트 정보 (report_config.md 없으므로 보고서에서 추론) ──
PROJECT_NAME = "서울 아파트 실거래가 분석"
PROJECT_SUBTITLE = "거래금액 결정 요인 통계 분석"
BACKGROUND = "서울 부동산 시장의 가격 변동성이 커지면서, 아파트 거래금액에 영향을 미치는 요인을 데이터 기반으로 이해할 필요성이 높아졌습니다."
PURPOSE = "서울 아파트 실거래 데이터(2020~2024)를 분석하여 거래금액에 영향을 미치는 핵심 요인을 파악합니다."
QUESTION = "어떤 요인(지역, 면적, 단지 규모, 노후도 등)이 서울 아파트 거래금액에 얼마나 영향을 미치는가?"
DATA_SOURCE = "국토교통부 실거래가 공개시스템 (APT_Data_Prep.csv)"
PRESENTER = ""
DATE = "2026-03-28"


# ── 헬퍼 함수 ──

def new_slide(prs):
    """빈 슬라이드 추가"""
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_title(slide, text, top=Inches(0.4), font_size=Pt(28)):
    """슬라이드 상단 제목"""
    txBox = slide.shapes.add_textbox(Inches(0.6), top, Inches(12.1), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.name = FONT
    p.font.color.rgb = CLR_TITLE
    p.font.bold = True

    # 제목 하단 구분선 (얇은 사각형)
    line = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0.6), top + Inches(0.75),
        Inches(12.1), Pt(2)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = CLR_ACCENT
    line.line.fill.background()


def add_bullet(tf, text, level=0, font_size=Pt(15), color=None):
    """글머리 항목 추가"""
    p = tf.add_paragraph()
    p.text = text
    p.level = level
    p.font.size = font_size
    p.font.name = FONT
    p.font.color.rgb = color if color else CLR_BODY
    p.line_spacing = Pt(22)


# ── 슬라이드 생성 함수 ──

def slide_cover(prs):
    """1. 표지"""
    slide = new_slide(prs)

    # 배경 직사각형 (상단 40%)
    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(3.8))
    bg.fill.solid()
    bg.fill.fore_color.rgb = CLR_COVER_BG
    bg.line.fill.background()

    # 프로젝트명
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.9), Inches(11.7), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = PROJECT_NAME
    p.font.size = Pt(40)
    p.font.name = FONT
    p.font.color.rgb = CLR_WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # 부제
    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.7), Inches(0.8))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = PROJECT_SUBTITLE
    p2.font.size = Pt(22)
    p2.font.name = FONT
    p2.font.color.rgb = RGBColor(0xAA, 0xCC, 0xEE)
    p2.alignment = PP_ALIGN.CENTER

    # 데이터 출처 + 날짜 (하단)
    txBox3 = slide.shapes.add_textbox(Inches(0.8), Inches(4.3), Inches(11.7), Inches(1.2))
    tf3 = txBox3.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    p3.text = f"데이터: {DATA_SOURCE}"
    p3.font.size = Pt(16)
    p3.font.name = FONT
    p3.font.color.rgb = CLR_BODY
    p3.alignment = PP_ALIGN.CENTER

    p4 = tf3.add_paragraph()
    p4.text = DATE
    p4.font.size = Pt(15)
    p4.font.name = FONT
    p4.font.color.rgb = CLR_BODY
    p4.alignment = PP_ALIGN.CENTER


def slide_problem(prs):
    """2. 문제 정의"""
    slide = new_slide(prs)
    add_title(slide, "문제 정의")

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(12.1), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "🔍  배경"
    p.font.size = Pt(18)
    p.font.name = FONT
    p.font.color.rgb = CLR_ACCENT
    p.font.bold = True

    add_bullet(tf, BACKGROUND, font_size=Pt(15))
    add_bullet(tf, "", font_size=Pt(8))

    add_bullet(tf, "🎯  분석 목적", font_size=Pt(18), color=CLR_ACCENT)
    # 마지막 단락 볼드 처리
    p2 = tf.paragraphs[-1]
    p2.font.bold = True

    add_bullet(tf, PURPOSE, font_size=Pt(15))
    add_bullet(tf, "", font_size=Pt(8))

    add_bullet(tf, "❓  핵심 질문", font_size=Pt(18), color=CLR_ACCENT)
    p3 = tf.paragraphs[-1]
    p3.font.bold = True

    add_bullet(tf, QUESTION, font_size=Pt(15))


def slide_data_overview(prs):
    """3. 데이터 소개"""
    slide = new_slide(prs)
    add_title(slide, "데이터 소개")

    # 왼쪽: 데이터 개요 테이블
    headers = ["항목", "값"]
    data = [
        ("데이터 출처", "국토교통부 실거래가 (APT_Data_Prep.csv)"),
        ("수집 기간", "2020-01 ~ 2024-12 (5개년)"),
        ("원본 규모", "220,890행 × 17열"),
        ("분석 데이터", "215,144행 × 14열 (전처리 후)"),
        ("수치형 변수", "10개 (거래금액, 전용면적, 층 등)"),
        ("범주형 변수", "7개 (시군구, 재건축 등)"),
        ("종속변수", "거래금액 (억 원)"),
    ]

    rows, cols = len(data) + 1, 2
    table = slide.shapes.add_table(
        rows, cols,
        Inches(0.6), Inches(1.5),
        Inches(12.1), Inches(0.42 * rows + 0.1)
    ).table
    table.columns[0].width = Inches(3.5)
    table.columns[1].width = Inches(8.6)

    # 헤더
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = CLR_TBL_HEADER
        for p in cell.text_frame.paragraphs:
            p.font.color.rgb = CLR_WHITE
            p.font.size = Pt(13)
            p.font.name = FONT
            p.font.bold = True

    # 데이터
    for r, (k, v) in enumerate(data, start=1):
        for c, val in enumerate([k, v]):
            cell = table.cell(r, c)
            cell.text = val
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = CLR_LIGHT_GRAY
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.name = FONT
                p.font.color.rgb = CLR_BODY


def slide_data_quality(prs):
    """4. 데이터 품질"""
    slide = new_slide(prs)
    add_title(slide, "데이터 품질")

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(12.1), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "✅  결측값 없음 — 전체 17개 컬럼 결측 0건"
    p.font.size = Pt(16)
    p.font.name = FONT
    p.font.color.rgb = RGBColor(0x27, 0xAE, 0x60)
    p.font.bold = True

    add_bullet(tf, "", font_size=Pt(8))
    add_bullet(tf, "⚠️  이상치 현황", font_size=Pt(17), color=CLR_ALERT)
    p2 = tf.paragraphs[-1]
    p2.font.bold = True

    add_bullet(tf, "층  min = -3  →  반지하/지하층 (도메인상 정상, 52건 유지)", font_size=Pt(14))
    add_bullet(tf, "경과년수  min = -2  →  입주 전 분양권 전매 (29건, 0으로 클리핑)", font_size=Pt(14))
    add_bullet(tf, "거래금액  평균 10.1억 > 중위수 8.3억  →  우측 편향 (고가 아파트, 도메인 정상)", font_size=Pt(14))

    add_bullet(tf, "", font_size=Pt(8))
    add_bullet(tf, "🔧  전처리 요약", font_size=Pt(17), color=CLR_ACCENT)
    p3 = tf.paragraphs[-1]
    p3.font.bold = True

    add_bullet(tf, "컬럼 제거: 시도명(단일값), 입주년도(경과년수와 r=-0.99), 주차대수(세대수와 r=0.95)", font_size=Pt(14))
    add_bullet(tf, "중복 제거: 5,746행 (-2.6%) → 최종 215,144행", font_size=Pt(14))
    add_bullet(tf, "타입 변환: 계약일자 → datetime64", font_size=Pt(14))


def slide_preprocessing(prs):
    """5. 전처리 과정"""
    slide = new_slide(prs)
    add_title(slide, "전처리 과정")

    # 전후 비교 테이블
    headers = ["항목", "전처리 전", "전처리 후", "변화"]
    data = [
        ("행 수", "220,890행", "215,144행", "-5,746행 (-2.6%)"),
        ("열 수", "17열", "14열", "-3열"),
        ("결측값", "0건", "0건", "—"),
        ("중복 행", "5,746행", "0행", "제거 완료"),
        ("다중공선성", "입주년도↔경과년수 r=-0.99\n세대수↔주차대수 r=0.95", "해소", "2개 컬럼 제거"),
    ]

    rows, cols = len(data) + 1, 4
    col_widths = [Inches(2.2), Inches(3.2), Inches(2.5), Inches(4.2)]
    tbl_width = sum(w.inches for w in col_widths)

    table = slide.shapes.add_table(
        rows, cols,
        Inches(0.6), Inches(1.45),
        Inches(tbl_width), Inches(0.43 * rows + 0.1)
    ).table

    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = CLR_TBL_HEADER
        for p in cell.text_frame.paragraphs:
            p.font.color.rgb = CLR_WHITE
            p.font.size = Pt(12)
            p.font.name = FONT
            p.font.bold = True

    for r, row_data in enumerate(data, start=1):
        for c, val in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = val
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = CLR_LIGHT_GRAY
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.name = FONT
                p.font.color.rgb = CLR_BODY

    # 전처리 원칙 메모
    note = slide.shapes.add_textbox(Inches(0.6), Inches(5.7), Inches(12.1), Inches(1.5))
    tf = note.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "전처리 원칙: 정보 손실 최소화 — 불필요 컬럼 제거 → 이상치 도메인 판단 → 타입 변환 → 중복 제거"
    p.font.size = Pt(13)
    p.font.name = FONT
    p.font.color.rgb = CLR_BODY
    p.font.italic = True


def slide_stat_results(prs):
    """6. 통계 검정 결과"""
    slide = new_slide(prs)
    add_title(slide, "통계 검정 결과 (8개 가설 검정)")

    headers = ["#", "분석", "검정", "통계량", "p값", "효과 크기", "해석"]
    data = [
        ("1", "시군구별 거래금액", "Kruskal-Wallis", "85,790.08", "<0.0001", "η²=0.3931", "★★★ 매우 큰 효과"),
        ("2", "재건축별 거래금액", "Mann-Whitney U", "3.25×10⁹", "0.0003", "d=0.1443", "검정력 과잉 가능"),
        ("3", "전용면적↔거래금액", "Spearman", "ρ=0.6098", "<0.0001", "r²=0.3718", "★★★ 강한 상관"),
        ("4", "계약년도별 거래금액", "Kruskal-Wallis", "9,083.75", "<0.0001", "η²=0.0332", "★ 작은 효과"),
        ("5", "계약월별 거래금액", "Kruskal-Wallis", "3,180.25", "<0.0001", "η²=0.0089", "검정력 과잉 가능"),
        ("6", "층↔거래금액", "Spearman", "ρ=0.1674", "<0.0001", "r²=0.0280", "약한 상관"),
        ("7", "경과년수↔거래금액", "Spearman", "ρ=-0.1289", "<0.0001", "r²=0.0166", "약한 음의 상관"),
        ("8", "세대수↔거래금액", "Spearman", "ρ=0.3237", "<0.0001", "r²=0.1048", "★★ 중간 상관"),
    ]

    rows, cols = len(data) + 1, len(headers)
    col_widths = [Inches(0.4), Inches(2.4), Inches(2.0), Inches(1.6), Inches(1.1), Inches(1.4), Inches(3.0)]

    table = slide.shapes.add_table(
        rows, cols,
        Inches(0.3), Inches(1.45),
        Inches(12.7), Inches(0.38 * rows + 0.1)
    ).table

    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = CLR_TBL_HEADER
        for p in cell.text_frame.paragraphs:
            p.font.color.rgb = CLR_WHITE
            p.font.size = Pt(11)
            p.font.name = FONT
            p.font.bold = True

    for r, row_data in enumerate(data, start=1):
        for c, val in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = val
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = CLR_LIGHT_GRAY
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.name = FONT
                # 효과 큰 행 강조
                if r in [1, 3] and c == 6:
                    p.font.color.rgb = CLR_ACCENT
                    p.font.bold = True
                else:
                    p.font.color.rgb = CLR_BODY

    # 주의사항
    note = slide.shapes.add_textbox(Inches(0.3), Inches(5.5), Inches(12.7), Inches(1.7))
    tf = note.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "⚠️  대표본 주의 (n=215,144): 8개 검정 모두 H0 기각이나, 재건축·계약월·층·경과년수는 효과 크기가 매우 작아 실질적 의미가 미미합니다."
    p.font.size = Pt(12)
    p.font.name = FONT
    p.font.color.rgb = CLR_ALERT
    p.font.italic = True


def slide_key_findings(prs):
    """7. 핵심 발견"""
    slide = new_slide(prs)
    add_title(slide, "핵심 발견")

    findings = [
        ("1위", "시군구", "η²=0.39", "거래금액의 39%를 설명하는 가장 강력한 요인.\n강남·서초·용산 등 고가 지역과 노원·도봉·강북 등 저가 지역의 격차가 뚜렷하다.\n25개 구 중 91.7% 쌍에서 유의한 차이 확인 (사후분석 Nemenyi)."),
        ("2위", "전용면적", "r²=0.37",  "거래금액과 강한 양의 상관관계 (ρ=0.61).\n면적이 넓을수록 가격이 높다는 직관적 결과.\n비선형 관계 가능성 → 추가 시각화 확인 권장."),
        ("3위", "세대수",  "r²=0.10",  "대단지일수록 가격이 높은 중간 수준의 상관 (ρ=0.32).\n대단지가 주로 강남권·신도시에 위치하는 입지 효과 반영 가능성."),
    ]

    colors = [CLR_ACCENT, RGBColor(0x27, 0xAE, 0x60), RGBColor(0x8E, 0x44, 0xAD)]

    y_start = Inches(1.4)
    box_h = Inches(1.7)
    gap = Inches(0.1)

    for idx, (rank, var, effect, desc) in enumerate(findings):
        y = y_start + idx * (box_h + gap)

        # 배경 박스
        bg = slide.shapes.add_shape(1, Inches(0.5), y, Inches(12.3), box_h)
        bg.fill.solid()
        bg.fill.fore_color.rgb = CLR_LIGHT_GRAY
        bg.line.fill.background()

        # 순위 + 변수명
        header_box = slide.shapes.add_textbox(Inches(0.7), y + Inches(0.12), Inches(3.5), Inches(0.55))
        tf = header_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{rank}  {var}  ({effect})"
        p.font.size = Pt(17)
        p.font.name = FONT
        p.font.color.rgb = colors[idx]
        p.font.bold = True

        # 설명
        desc_box = slide.shapes.add_textbox(Inches(0.7), y + Inches(0.65), Inches(12.0), Inches(1.0))
        tf2 = desc_box.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = desc
        p2.font.size = Pt(13)
        p2.font.name = FONT
        p2.font.color.rgb = CLR_BODY
        p2.line_spacing = Pt(20)


def slide_conclusion(prs):
    """8. 결론"""
    slide = new_slide(prs)
    add_title(slide, "결론")

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(12.1), Inches(5.8))
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "핵심 답변: 서울 아파트 거래금액은 지역(시군구)과 면적(전용면적)이 가장 크게 결정한다."
    p.font.size = Pt(18)
    p.font.name = FONT
    p.font.color.rgb = CLR_ACCENT
    p.font.bold = True

    add_bullet(tf, "", font_size=Pt(8))
    add_bullet(tf, "주요 결과", font_size=Pt(16), color=CLR_TITLE)
    tf.paragraphs[-1].font.bold = True

    add_bullet(tf, "시군구: η²=0.39 — 거래금액 분산의 39% 설명, 입지가 가격의 핵심 변수", font_size=Pt(14))
    add_bullet(tf, "전용면적: r²=0.37 — 면적이 넓을수록 가격 상승, 강한 양의 상관", font_size=Pt(14))
    add_bullet(tf, "재건축·계약월·층·경과년수: 통계적 유의성 있으나 효과 크기 매우 작음 (대표본 검정력 과잉)", font_size=Pt(14))

    add_bullet(tf, "", font_size=Pt(8))
    add_bullet(tf, "한계점", font_size=Pt(16), color=CLR_ALERT)
    tf.paragraphs[-1].font.bold = True

    add_bullet(tf, "비모수 검정 기반 — 그룹 간 순위 차이이며 평균 차이는 별도 확인 필요", font_size=Pt(13))
    add_bullet(tf, "거래금액 우편향 (평균 10.1억 > 중위수 8.3억) — 로그 변환 후 재분석 검토", font_size=Pt(13))
    add_bullet(tf, "시군구 효과가 다른 변수 효과를 혼재시킬 가능성 (교란 변수)", font_size=Pt(13))

    add_bullet(tf, "", font_size=Pt(8))
    add_bullet(tf, "후속 분석 제안", font_size=Pt(16), color=RGBColor(0x27, 0xAE, 0x60))
    tf.paragraphs[-1].font.bold = True

    add_bullet(tf, "시군구별 가격 분포 박스플롯 시각화 (/da-viz box 시군구 거래금액)", font_size=Pt(13))
    add_bullet(tf, "전용면적↔거래금액 비선형 관계 산점도 확인", font_size=Pt(13))
    add_bullet(tf, "다중회귀분석 — 시군구·전용면적·세대수 독립적 영향력 분리", font_size=Pt(13))


# ── 메인 ──

def build_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    print("슬라이드 생성 중...")
    slide_cover(prs)
    print("  [1/8] 표지")
    slide_problem(prs)
    print("  [2/8] 문제 정의")
    slide_data_overview(prs)
    print("  [3/8] 데이터 소개")
    slide_data_quality(prs)
    print("  [4/8] 데이터 품질")
    slide_preprocessing(prs)
    print("  [5/8] 전처리 과정")
    slide_stat_results(prs)
    print("  [6/8] 통계 검정 결과")
    slide_key_findings(prs)
    print("  [7/8] 핵심 발견")
    slide_conclusion(prs)
    print("  [8/8] 결론")

    OUTPUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PPTX))
    print(f"\nPPTX 저장 완료: {OUTPUT_PPTX}")


if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("--format", default="pptx", choices=["pptx", "html", "all"])
    args = parser.parse_args()

    if args.format in ("pptx", "all"):
        build_pptx()
